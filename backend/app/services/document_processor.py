"""
Universal document processor with full OCR fallback.

PDF parsing: pdfplumber (wraps pdfminer.six) is used first.
pdfminer.six has a known class of KeyError bugs on non-standard PDFs:
  - KeyError: '_type'   (PDFStream missing internal type key)
  - KeyError: 'N'       (colour space dict missing N key)
  - KeyError: 'Type'    (stream dict missing Type key)
  - KeyError: 'MediaBox'(page missing MediaBox)
  - TypeError:          (PDFObjRef not iterable in decode())
References:
  https://github.com/pdfminer/pdfminer.six/issues/559
  https://github.com/pdfminer/pdfminer.six/issues/551
  https://github.com/jsvine/pdfplumber/issues/1204

All of these are caught per-page AND at the document level, with
automatic fallback to pdf2image + pytesseract OCR.
"""
import json
import logging
import os
from pathlib import Path
from typing import List
import re

from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter

logger = logging.getLogger(__name__)
OCR_LANG = os.getenv("OCR_LANG", "eng")

SUPPORTED_EXTENSIONS = {
    "pdf", "docx", "doc", "xlsx", "xls",
    "pptx", "ppt", "csv", "txt", "md",
    "html", "htm", "json",
    "png", "jpg", "jpeg", "tiff", "tif", "bmp", "gif", "webp",
}


def _fix_ocr_spacing(text: str) -> str:
    """Fix isolated single chars injected by OCR mid-word e.g. 'Instruct i ons' -> 'Instructions'."""
    # Pattern 1: single lowercase char surrounded by spaces between word chars
    text = re.sub(r"(?<=[a-zA-Z]) ([a-z]) (?=[a-z])", r"\1", text)
    # Pattern 2: capitalised word fragment + single char + word fragment
    text = re.sub(r"\b([A-Z][a-z]+) ([a-z]) ([a-z]+)\b", r"\1\2\3", text)
    return text
def _score_text(text: str) -> int:
    # Prefer outputs with useful alphanumeric content, not just whitespace/noise.
    return sum(ch.isalnum() for ch in text)


def _ocr_variants(img):
    from PIL import Image, ImageEnhance, ImageFilter, ImageOps

    # Generate a small set of OCR-friendly variants; pick the best OCR result.
    base = ImageOps.exif_transpose(img).convert("RGB")

    # Upscale small images to improve character legibility.
    max_side = max(base.size)
    if max_side < 1400:
        scale = max(1.0, 1600 / max_side)
        base = base.resize(
            (int(base.width * scale), int(base.height * scale)),
            Image.Resampling.LANCZOS,
        )

    gray = ImageOps.grayscale(base)
    auto = ImageOps.autocontrast(gray)
    sharp = auto.filter(ImageFilter.SHARPEN)

    # Light binary threshold often helps low-contrast scans/screenshots.
    bw = sharp.point(lambda p: 255 if p > 170 else 0)

    # Slight contrast boost for faded pages.
    boosted = ImageEnhance.Contrast(auto).enhance(1.4)

    return [base, auto, sharp, bw, boosted]


def _ocr_image_best_effort(img) -> str:
    import pytesseract

    best_text = ""
    for variant in _ocr_variants(img):
        for cfg in (
            "--oem 1 --psm 6",
            "--oem 1 --psm 11",
            "--oem 1 --psm 3",
        ):
            try:
                txt = pytesseract.image_to_string(variant, lang=OCR_LANG, config=cfg) or ""
                if _score_text(txt) > _score_text(best_text):
                    best_text = txt
            except Exception:
                continue

    return best_text.strip()


async def process_file(
    file_path: str,
    original_filename: str,
    metadata: dict = None,
    chunk_size: int = 800,
    chunk_overlap: int = 120,
) -> List[Document]:
    """Process any supported file into LangChain Document chunks."""
    ext = Path(original_filename).suffix.lower().lstrip(".")

    # Build metadata with ONLY ChromaDB-safe primitives (str/int/float/bool).
    # None values and complex types are dropped — ChromaDB 0.5.x rejects them.
    base_meta: dict = {"source": str(original_filename), "file_type": str(ext)}
    if metadata:
        for k, v in metadata.items():
            if isinstance(v, (str, int, float, bool)):
                base_meta[k] = v
            elif v is not None:
                base_meta[k] = str(v)

    try:
        if ext == "pdf":
            content = _extract_pdf(file_path)
        elif ext in ("docx", "doc"):
            content = _extract_docx(file_path)
        elif ext in ("xlsx", "xls"):
            content = _extract_excel(file_path)
        elif ext in ("pptx", "ppt"):
            content = _extract_pptx(file_path)
        elif ext == "csv":
            content = _extract_csv(file_path)
        elif ext in ("html", "htm"):
            content = _extract_html(file_path)
        elif ext == "json":
            content = _extract_json(file_path)
        elif ext in ("png", "jpg", "jpeg", "tiff", "tif", "bmp", "gif", "webp"):
            content = _extract_image_ocr(file_path)
        else:
            content = _extract_text(file_path)
    except Exception as e:
        logger.error(f"Extraction failed for '{original_filename}': {e}", exc_info=True)
        raise


    content = _fix_ocr_spacing(content)
    if not content or not content.strip():
        raise ValueError(
            f"No text could be extracted from '{original_filename}'. "
            "The file may be empty, image-only, or password-protected."
        )

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        separators=["\n\n", "\n", ". ", " ", ""],
    )
    docs = splitter.create_documents([content], metadatas=[base_meta])
    logger.info(f"'{original_filename}' → {len(docs)} chunks (ext={ext})")
    kb_id = (metadata or {}).get("kb_id")
    if kb_id is not None:
        try:
            from .kb_manifest import upsert_manifest
            upsert_manifest(int(kb_id), original_filename, content, ext)
        except Exception as _me:
            logger.warning("manifest upsert: %s", _me)
    return docs


# ── PDF extractor (most complex — full pdfminer error shielding) ──────────────

def _extract_pdf(file_path: str) -> str:
    """
    Strategy:
      1. pdfplumber page-by-page — catch pdfminer KeyError/TypeError per page
         and per document (these are upstream pdfminer.six bugs on non-standard PDFs)
      2. If total extracted text < 50 chars OR pdfplumber raised at document level
         → fall back to pdf2image + pytesseract OCR
         (pdf2image uses poppler to render pages, completely bypasses pdfminer)
    """
    import pdfplumber

    text = ""
    pdfplumber_doc_error = None

    # ── Attempt 1: pdfplumber ─────────────────────────────────────────────────
    try:
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages, start=1):
                try:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
                except (KeyError, TypeError, AttributeError, Exception) as page_err:
                    # pdfminer raises KeyError for missing stream attributes:
                    #   '_type', 'N', 'Type', 'MediaBox', etc.
                    # TypeError raised when PDFObjRef encountered in decode()
                    logger.warning(
                        f"pdfplumber: page {page_num} failed "
                        f"({type(page_err).__name__}: {page_err!r}) — "
                        f"this page will be covered by OCR fallback"
                    )
                    # Mark text as insufficient to force OCR
                    text = ""
                    break

    except (KeyError, TypeError, AttributeError, Exception) as doc_err:
        # Entire document-level open/parse failed
        pdfplumber_doc_error = doc_err
        logger.warning(
            f"pdfplumber failed at document level for '{file_path}' "
            f"({type(doc_err).__name__}: {doc_err!r}) — using full OCR"
        )
        text = ""

    # ── Attempt 2: OCR fallback via poppler + tesseract ───────────────────────
    # Triggered when: pdfplumber failed entirely, OR yielded < 50 chars
    if len(text.strip()) < 50:
        logger.info(
            f"OCR fallback triggered for '{file_path}' "
            f"(pdfplumber yielded {len(text.strip())} chars)"
        )
        try:
            from pdf2image import convert_from_path

            images = convert_from_path(file_path, dpi=300)
            ocr_text = ""
            for img in images:
                page_ocr = _ocr_image_best_effort(img)
                if page_ocr:
                    ocr_text += page_ocr + "\n\n"

            if ocr_text.strip():
                logger.info(
                    f"OCR succeeded: {len(ocr_text.strip())} chars from '{file_path}'"
                )
                return ocr_text.strip()

            # OCR ran but produced nothing (e.g. blank pages, pure graphics)
            raise ValueError(
                "OCR found no readable text in this PDF. "
                "The file may contain only images/graphics with no text."
            )

        except (ImportError, Exception) as ocr_err:
            logger.error(f"OCR fallback failed: {ocr_err}", exc_info=True)
            # If we have any partial pdfplumber text, return it
            if text.strip():
                logger.warning("Returning partial pdfplumber text as last resort")
                return text.strip()
            # Nothing worked — raise descriptive error
            primary = f"pdfminer: {pdfplumber_doc_error!r}" if pdfplumber_doc_error else "low text yield"
            raise ValueError(
                f"Could not extract text from PDF '{Path(file_path).name}'. "
                f"Primary: {primary}. OCR: {ocr_err}. "
                "Try re-saving as PDF (File → Print → Save as PDF) and re-uploading."
            ) from ocr_err

    return text.strip()


# ── Other extractors ──────────────────────────────────────────────────────────

def _extract_docx(file_path: str) -> str:
    from docx import Document as DocxDoc

    doc = DocxDoc(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def _extract_excel(file_path: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
    parts = []
    for sheet_name in wb.sheetnames:
        parts.append(f"=== Sheet: {sheet_name} ===")
        sheet = wb[sheet_name]
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(v) for v in row if v is not None)
            if row_text.strip():
                parts.append(row_text)
    return "\n".join(parts)


def _extract_pptx(file_path: str) -> str:
    from pptx import Presentation

    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- Slide {i} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
    return "\n\n".join(parts)


def _extract_csv(file_path: str) -> str:
    import pandas as pd

    df = pd.read_csv(file_path)
    return df.to_string(index=False)


def _extract_html(file_path: str) -> str:
    from bs4 import BeautifulSoup

    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "html.parser")
    return soup.get_text(separator="\n", strip=True)


def _extract_json(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8") as f:
        data = json.load(f)
    return json.dumps(data, indent=2, ensure_ascii=False)


def _extract_image_ocr(file_path: str) -> str:
    from PIL import Image

    img = Image.open(file_path)
    return _ocr_image_best_effort(img)


def _extract_text(file_path: str) -> str:
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()
    
